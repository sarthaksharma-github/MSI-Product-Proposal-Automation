# ============================================================
#  MSI SERVICES — SLIDE AUTOMATION TOOL v3.5
#  Streamlit App
# ============================================================

import streamlit as st
import pandas as pd
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image as PILImage
import io, copy, os, traceback, re, json, base64
import streamlit.components.v1 as components

# ══════════════════════════════════════════════════════════════
#  MSI LOGO AUTO-COPY ROUTINE
# ══════════════════════════════════════════════════════════════
import shutil
project_dir = os.path.dirname(os.path.abspath(__file__))
dst_logo = os.path.join(project_dir, "msi_logo.png")

# Try to find logo in brain folder or current dir
logo_copied = False
if not os.path.exists(dst_logo):
    brain_dir = r"C:\Users\sarthak.s\.gemini\antigravity\brain\c86fcd1e-6a81-44ad-9eeb-5fce89a06592"
    possible_logos = ["media__1782309457101.png", "media__1782309064522.png"]
    for l_name in possible_logos:
        src_path = os.path.join(brain_dir, l_name)
        if os.path.exists(src_path):
            try:
                shutil.copy(src_path, dst_logo)
                logo_copied = True
                break
            except:
                pass

# ══════════════════════════════════════════════════════════════
#  PAGE CONFIG
# ══════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="MSI Slide Automation Tool",
    page_icon="📊",
    layout="wide"
)

# ══════════════════════════════════════════════════════════════
#  PREMIUM STYLING
# ══════════════════════════════════════════════════════════════

st.markdown("""
<style>
  @import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@300;400;500;600;700;800&display=swap');

  /* Global resets and body styling */
  html, body, [class*="css"], [class*="st-"] {
    font-family: 'Plus Jakarta Sans', sans-serif !important;
    background-color: #FAF8F5 !important;
  }
  
  .main .block-container { 
    padding-top: 1.5rem !important; 
    max-width: 1250px !important; 
  }
  
  #MainMenu, footer, header { visibility: hidden; }

  /* Style all vertical block border containers (cards) */
  div[data-testid="stVerticalBlockBorder"] {
    background-color: #FFFFFF !important;
    border: 1.5px solid #EDE8E1 !important;
    border-radius: 12px !important;
    padding: 24px !important;
    box-shadow: 0 2px 12px rgba(139,111,78,0.03) !important;
    margin-bottom: 16px !important;
  }

  /* Left Sidebar Styling */
  [data-testid="stSidebar"] {
    background-color: #FFFFFF !important;
    border-right: 1px solid #EDE8E1 !important;
  }
  [data-testid="stSidebarNav"] {
    display: none !important;
  }
  
  /* Sidebar label styling */
  .menu-label {
    font-size: 10px;
    font-weight: 700;
    color: #B8A898;
    letter-spacing: 1.2px;
    margin-bottom: 8px;
    padding: 0 12px;
  }

  /* Style Streamlit sidebar buttons to look like nav links */
  div[data-testid="stSidebar"] button {
    background-color: transparent !important;
    color: #5C483A !important;
    border: 1px solid transparent !important;
    text-align: left !important;
    justify-content: flex-start !important;
    font-size: 13px !important;
    font-weight: 600 !important;
    border-radius: 8px !important;
    padding: 10px 12px !important;
    margin-bottom: 6px !important;
    transition: all 0.2s ease !important;
    width: 100% !important;
  }
  div[data-testid="stSidebar"] button:hover {
    background-color: #FFF8F2 !important;
    color: #8B6F4E !important;
    border-color: #EDE8E1 !important;
  }
  /* Active button styling */
  div[data-testid="stSidebar"] button[kind="primary"] {
    background-color: #FFF3E8 !important;
    color: #8B6F4E !important;
    border-left: 3px solid #8B6F4E !important;
    border-radius: 0 8px 8px 0 !important;
    font-weight: 700 !important;
    box-shadow: none !important;
  }

  /* Step card numbers */
  .step-card-num {
    width: 32px;
    height: 32px;
    border-radius: 50%;
    background: #8B6F4E;
    color: #ffffff;
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 13px;
    font-weight: 700;
    flex-shrink: 0;
  }

  /* File upload dropzone enhancement */
  [data-testid="stFileUploader"] > div {
    border: 2px dashed #CFC0B0 !important;
    border-radius: 12px !important;
    background: #FAF8F5 !important;
    transition: border-color 0.3s, background 0.3s, box-shadow 0.3s !important;
  }
  [data-testid="stFileUploader"] > div:hover {
    border-color: #8B6F4E !important;
    background: #FFF8F2 !important;
    box-shadow: 0 0 0 4px rgba(139,111,78,0.1) !important;
  }

  /* Text inputs and number inputs */
  div[data-testid="stTextInput"] input,
  div[data-testid="stNumberInput"] input {
    border-radius: 8px !important;
    border: 1.5px solid #E0D5CA !important;
    background-color: #FFFFFF !important;
    transition: border-color 0.25s, box-shadow 0.25s !important;
  }
  div[data-testid="stTextInput"] input:focus,
  div[data-testid="stNumberInput"] input:focus {
    border-color: #8B6F4E !important;
    box-shadow: 0 0 0 3px rgba(139,111,78,0.12) !important;
  }

  /* Primary conversion action buttons */
  .stButton > button[kind="primary"], .stDownloadButton > button {
    background: linear-gradient(135deg, #8B6F4E 0%, #7A6042 100%) !important;
    color: #ffffff !important;
    font-weight: 700 !important;
    font-size: 13px !important;
    letter-spacing: 0.3px !important;
    border: none !important;
    border-radius: 8px !important;
    padding: 10px 24px !important;
    width: 100% !important;
    box-shadow: 0 4px 12px rgba(139,111,78,0.25) !important;
    transition: all 0.2s cubic-bezier(.4,0,.2,1) !important;
  }
  .stButton > button[kind="primary"]:hover, .stDownloadButton > button:hover {
    background: linear-gradient(135deg, #7A6042 0%, #6A5235 100%) !important;
    transform: translateY(-1px) !important;
    box-shadow: 0 6px 16px rgba(139,111,78,0.3) !important;
  }

  /* Stepper timeline on the right */
  .timeline-container {
    padding: 0;
  }
  .timeline-title {
    font-size: 13px;
    font-weight: 700;
    color: #2C1F14;
    margin-bottom: 16px;
    text-transform: uppercase;
    letter-spacing: 0.5px;
  }
  .timeline-item {
    display: flex;
    gap: 12px;
    position: relative;
    padding-bottom: 20px;
  }
  .timeline-item:not(:last-child)::after {
    content: '';
    position: absolute;
    left: 13px;
    top: 26px;
    width: 2px;
    height: calc(100% - 24px);
    background-color: #EDE8E1;
  }
  .timeline-item.done:not(:last-child)::after {
    background-color: #8B6F4E;
  }
  .timeline-badge {
    width: 26px;
    height: 26px;
    border-radius: 50%;
    background: #FAF8F5;
    border: 2px solid #EDE8E1;
    color: #B8A898;
    display: flex;
    align-items: center;
    justify-content: center;
    font-size: 11px;
    font-weight: 700;
    z-index: 1;
    transition: all 0.3s;
  }
  .timeline-item.done .timeline-badge {
    background: #8B6F4E;
    border-color: #8B6F4E;
    color: #FFFFFF;
  }
  .timeline-item.active .timeline-badge {
    background: #FFF3E8;
    border-color: #8B6F4E;
    color: #8B6F4E;
    box-shadow: 0 0 0 4px rgba(139,111,78,0.12);
  }
  .timeline-content {
    flex: 1;
  }
  .timeline-label {
    font-size: 12px;
    font-weight: 700;
    color: #2C1F14;
  }
  .timeline-desc {
    font-size: 11px;
    color: #9A8070;
    margin-top: 1px;
    word-break: break-all;
  }

  /* Success result card */
  .result-card {
    background: linear-gradient(135deg, #FFF8F2 0%, #FFF3E8 100%);
    border: 1.5px solid #D4B896;
    border-radius: 12px;
    padding: 20px;
    text-align: center;
  }
  .result-icon { font-size: 36px; line-height: 1; margin-bottom: 8px; }
  .result-title { font-size: 16px; font-weight: 800; color: #2C1F14; margin: 0 0 2px; }
  .result-sub { font-size: 11.5px; color: #8B6F4E; font-weight: 500; margin: 0 0 12px; word-break: break-all; }
  .result-stats {
    display: flex; justify-content: center; gap: 12px; flex-wrap: wrap;
  }
  .stat-chip {
    background: #ffffff; border: 1px solid #E8D9CA; border-radius: 8px;
    padding: 6px 10px; text-align: center; min-width: 65px;
  }
  .stat-chip .val { font-size: 15px; font-weight: 800; color: #8B6F4E; display: block; }
  .stat-chip .lbl { font-size: 8px; color: #9A8070; font-weight: 600; text-transform: uppercase; letter-spacing: 0.5px; }

  /* Dividers and Footers */
  .fancy-divider { height: 1px; background: linear-gradient(90deg, transparent, #E0D5CA, transparent); margin: 20px 0; }
  .msi-footer {
    border-top: 1px solid #EDE8E1;
    padding: 16px 0; font-size: 11px; color: #B8A898;
    text-align: center; margin-top: 40px;
    letter-spacing: 0.3px;
  }

  /* Popover styling overrides */
  div[data-testid="stPopoverBody"] {
    border: 1.5px solid #EDE8E1 !important;
    border-radius: 10px !important;
    box-shadow: 0 4px 16px rgba(139,111,78,0.1) !important;
  }
  div[data-testid="stPopover"] > button {
    background: #FFFFFF !important;
    border: 1.5px solid #EDE8E1 !important;
    color: #5C483A !important;
    font-weight: 600 !important;
    font-size: 13px !important;
    box-shadow: none !important;
    border-radius: 8px !important;
  }
  div[data-testid="stPopover"] > button:hover {
    background: #FFF8F2 !important;
    border-color: #8B6F4E !important;
    color: #8B6F4E !important;
  }
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
#  CORE ENGINE
# ══════════════════════════════════════════════════════════════

def clone_slide(prs, source_slide):
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    sp_tree = new_slide.shapes._spTree
    while len(sp_tree) > 2:
        sp_tree.remove(sp_tree[-1])
    source_sp_tree = source_slide.shapes._spTree
    for el in list(source_sp_tree)[2:]:
        sp_tree.append(copy.deepcopy(el))
    return new_slide

def get_excel_images(excel_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(excel_bytes))
    ws = wb.active
    raw = []
    for img in ws._images:
        try:
            row = img.anchor._from.row
            col = img.anchor._from.col
            raw.append((row, col, img._data()))
        except:
            pass
    raw.sort(key=lambda x: x[0])
    col_to_images = {}
    for row, col, img_bytes in raw:
        col_to_images.setdefault(col, []).append(img_bytes)
    return col_to_images

def parse_placeholder_tag(full_tag):
    inner = full_tag.strip().lstrip('[').rstrip(']').strip()
    m = re.match(r'^(.+?)\s*\(([^)]+)\)\s*$', inner)
    if m:
        base = m.group(1).strip()
        ann  = m.group(2).strip()
        ann_l = ann.lower()
        if any(x in ann_l for x in ['image', 'img', 'photo', 'picture']):
            return base, 'Image', None
        if 'currency' in ann_l or any(c in ann for c in '$€£¥₹'):
            sym_m = re.search(r'[\.\$€£¥₹]', ann)
            symbol = sym_m.group(0) if sym_m else '$'
            return base, 'Currency', symbol
        if any(x in ann_l for x in ['percent', '%']):
            return base, 'Percentage', None
        if any(x in ann_l for x in ['integer', 'int', 'number', 'num']):
            return base, 'Integer', None
        return base, 'Text', None
    return inner, None, None


def build_auto_mapping(all_pptx_tags, excel_columns):
    def _norm(s):
        s = re.sub(r'[\[\]_]', ' ', str(s))
        return re.sub(r'\s+', ' ', s).strip().lower()

    def _num(s):
        m = re.search(r'(\d+)\s*$', s.strip())
        return int(m.group(1)) if m else None

    def _best_excel_match(base_norm, cols):
        tag_num   = _num(base_norm)
        tag_roots = {w for w in base_norm.split() if not w.isdigit()}
        best_col, best_score = '', -999
        for col in cols:
            col_norm  = _norm(col)
            col_num   = _num(col_norm)
            col_roots = {w for w in col_norm.split() if not w.isdigit()}
            if not any(tr in cr or cr in tr for tr in tag_roots for cr in col_roots):
                continue
            score = 1
            if tag_num is not None and col_num is not None:
                score += 5 if tag_num == col_num else -10
            elif tag_num is not None:
                score -= 1
            if base_norm in col_norm or col_norm in base_norm:
                score += 2
            if score > best_score:
                best_score, best_col = score, col
        return best_col if best_score > 0 else ''

    mapping_dict   = {}
    image_mappings = {}

    for full_tag in all_pptx_tags:
        base, col_type, symbol = parse_placeholder_tag(full_tag)
        base_norm   = _norm(base)
        matched_col = _best_excel_match(base_norm, excel_columns)

        if not matched_col:
            continue

        if col_type == 'Image':
            image_mappings[full_tag] = matched_col
        else:
            if col_type in ('Currency', 'Percentage', 'Integer', 'Text'):
                fmt = col_type
            else:
                tag_l = full_tag.lower()
                if any(x in tag_l for x in ['retail', 'cost', 'price', 'rtl', 'amt', 'value']):
                    fmt    = 'Currency'
                    symbol = symbol or '$'
                elif any(x in tag_l for x in ['imu', 'percent', 'pct', '%']):
                    fmt = 'Percentage'
                elif any(x in tag_l for x in ['units', 'qty', 'count']):
                    fmt = 'Integer'
                else:
                    fmt = 'Text'

            mapping_dict[full_tag] = {
                'column': matched_col,
                'format': fmt,
                'symbol': symbol or '$',
            }

    return mapping_dict, image_mappings


def process_text_frame(tf, placeholders):
    for para in tf.paragraphs:
        for key, val in placeholders.items():
            if key in para.text:
                new_text = para.text.replace(key, str(val))
                if para.runs:
                    para.runs[0].text = new_text
                    for i in range(1, len(para.runs)):
                        para.runs[i].text = ""

def purge_empty_paragraphs(shape):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    txBody = tf._txBody
    all_paras = tf.paragraphs
    if len(all_paras) <= 1:
        return
    for para in list(all_paras):
        run_text = ''.join(r.text or '' for r in para.runs).strip()
        if run_text == '' and len(tf.paragraphs) > 1:
            txBody.remove(para._p)

def replace_text_in_shape(shape, placeholders):
    if shape.has_text_frame:
        process_text_frame(shape.text_frame, placeholders)
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                process_text_frame(cell.text_frame, placeholders)

def find_image_placeholder(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if "insert product" in text or "picture here" in text:
                return shape
    return None

def insert_image_into_placeholder(slide, img_bytes, placeholder):
    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
    placeholder.fill.background()
    for para in placeholder.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""

    img = PILImage.open(io.BytesIO(img_bytes)).convert('RGB')
    img_w, img_h = img.size
    target_ratio, img_ratio = width / height, img_w / img_h

    if img_ratio > target_ratio:
        new_h = int(img_w / target_ratio)
        padded = PILImage.new('RGB', (img_w, new_h), (255, 255, 255))
        padded.paste(img, (0, (new_h - img_h) // 2))
    else:
        new_w = int(img_h * target_ratio)
        padded = PILImage.new('RGB', (new_w, img_h), (255, 255, 255))
        padded.paste(img, (((new_w - img_w) // 2), 0))

    output = io.BytesIO()
    padded.save(output, format='JPEG', quality=95)
    output.seek(0)
    slide.shapes.add_picture(output, left, top, width, height)

def safe_format(val, is_currency=False, currency_symbol='$', is_percent=False):
    try:
        cleaned = re.sub(r'[^\d.\-]', '', str(val).strip())
        num = 0.0 if (not cleaned or cleaned.lower() == 'nan') else float(cleaned)
        if is_currency:
            return f"{currency_symbol}{num:,.2f}" if num % 1 != 0 else f"{currency_symbol}{num:,.0f}"
        if is_percent:
            return f"{num:.0%}"
        return f"{num:,.0f}"
    except:
        return str(val)

def safe_text(val):
    v = str(val).strip()
    return "" if v.lower() == 'nan' else v

def extract_placeholders_from_pptx(pptx_bytes):
    placeholders = set()
    pattern = re.compile(r'\[([^\]]+)\]')
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for match in pattern.findall(para.text):
                            placeholders.add(f"[{match.strip()}]")
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for match in pattern.findall(para.text):
                                    placeholders.add(f"[{match.strip()}]")
    except Exception as e:
        pass
    return sorted(list(placeholders))

CONFIG_FILE = "mapping_config.json"

def load_mapping_config():
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, 'r') as f:
                return json.load(f)
        except:
            pass
    return {}

def save_mapping_config(config):
    try:
        with open(CONFIG_FILE, 'w') as f:
            json.dump(config, f, indent=4)
    except:
        pass

def render_slide_preview(mapping_dict, image_mappings=None, excel_row=None, image_map=None, excel_row_idx=None, is_template_mode=False):
    if image_mappings is None:
        image_mappings = {}
    if image_map is None:
        image_map = {}
        
    title_tag = None
    sku_tag = None
    feature_tags = []
    benefit_tags = []
    metadata_tags = []
    
    for tag in mapping_dict.keys():
        tag_lower = tag.lower()
        if not title_tag and ("name" in tag_lower or "title" in tag_lower):
            title_tag = tag
        elif not sku_tag and ("num" in tag_lower or "sku" in tag_lower or "id" in tag_lower or "#" in tag_lower):
            sku_tag = tag
        elif "feat" in tag_lower:
            feature_tags.append(tag)
        elif "ben" in tag_lower:
            benefit_tags.append(tag)
        else:
            metadata_tags.append(tag)
            
    feature_tags.sort()
    benefit_tags.sort()
    metadata_tags.sort()
    
    all_tags = list(mapping_dict.keys())
    if not title_tag and all_tags:
        title_tag = all_tags[0]
        
    def get_val(tag):
        if is_template_mode:
            return tag
        col = mapping_dict[tag].get("column", "")
        fmt = mapping_dict[tag].get("format", "Text")
        if not col or excel_row is None:
            return f"({tag} Unmapped)"
        raw_val = excel_row.get(col, "")
        if pd.isna(raw_val):
            return ""
        if fmt == "Currency":
            return safe_format(raw_val, is_currency=True)
        elif fmt == "Percentage":
            return safe_format(raw_val, is_percent=True)
        elif fmt == "Integer":
            return safe_format(raw_val)
        return safe_text(raw_val)
        
    title_val = get_val(title_tag) if title_tag else "Sample Product Title"
    sku_val = get_val(sku_tag) if sku_tag else "SKU-000000"
    
    features_html = ""
    for f_tag in feature_tags:
        val = get_val(f_tag)
        if val:
            features_html += f"<li>{val}</li>"
    if not features_html:
        features_html = "<li>Key Product Feature #1</li><li>Key Product Feature #2</li>"
        
    benefits_html = ""
    for b_tag in benefit_tags:
        val = get_val(b_tag)
        if val:
            benefits_html += f"<li>{val}</li>"
    if not benefits_html:
        benefits_html = "<li>Customer Benefit #1</li><li>Customer Benefit #2</li>"

    details_html = ""
    for m_tag in metadata_tags:
        val = get_val(m_tag)
        label = m_tag.replace("[", "").replace("]", "").replace("_", " ").title()
        details_html += f"""
        <div style="display: flex; justify-content: space-between; padding: 4px 0; border-bottom: 1px solid #F0EAE1;">
            <span style="font-weight: 600; color: #5C483A; font-size: 12px;">{label}</span>
            <span style="color: #2C1F14; font-size: 12px; font-weight: 500;">{val or '—'}</span>
        </div>
        """
        
    def _img_box(label, img_bytes_val):
        if is_template_mode:
            return f'<div style="flex:1;min-width:140px;height:200px;background:#EFECE8;border:2px dashed #C8B8AA;border-radius:6px;display:flex;align-items:center;justify-content:center;color:#8B6F4E;font-weight:600;font-size:12px;text-align:center;padding:6px;">{label}</div>'
        if img_bytes_val:
            b64 = base64.b64encode(img_bytes_val).decode('utf-8')
            return f'<img src="data:image/jpeg;base64,{b64}" style="flex:1;min-width:140px;height:200px;object-fit:contain;background:#FFF;border-radius:6px;border:1px solid #E4D6CA;padding:4px;"/>'
        return f'<div style="flex:1;min-width:140px;height:200px;background:#F9F6F3;border:1px solid #E4D6CA;border-radius:6px;display:flex;align-items:center;justify-content:center;color:#C4B0A0;font-style:italic;font-size:12px;">No image for {label}</div>'

    images_html_parts = []
    if image_mappings:
        for img_tag, col_name in image_mappings.items():
            if is_template_mode:
                images_html_parts.append(_img_box(img_tag, None))
            else:
                col_idx = None
                if excel_row is not None and col_name:
                    cols = list(excel_row.keys())
                    if col_name in cols:
                        col_idx = cols.index(col_name)
                img_bytes_val = image_map.get((excel_row_idx, col_idx)) if col_idx is not None and excel_row_idx is not None else None
                images_html_parts.append(_img_box(img_tag, img_bytes_val))
    else:
        images_html_parts.append(_img_box('[IMAGE]', None))

    img_html = f'<div style="display:flex;gap:10px;flex-wrap:wrap;">{" ".join(images_html_parts)}</div>'
            
    html_content = f"""
    <div style="font-family: 'DM Sans', sans-serif; background: #FFFFFF; border: 1px solid #E4D6CA; border-radius: 10px; box-shadow: 0 4px 12px rgba(44, 31, 20, 0.06); padding: 20px; max-width: 100%; margin: 10px auto;">
        <div style="display: flex; justify-content: space-between; align-items: baseline; border-bottom: 2px solid #8B6F4E; padding-bottom: 8px; margin-bottom: 15px;">
            <h3 style="margin: 0; color: #2C1F14; font-size: 18px; font-weight: 600;">{title_val}</h3>
            <span style="font-size: 11px; background: #FFF3E8; border: 1px solid #E4D6CA; color: #8B6F4E; padding: 2px 8px; border-radius: 20px; font-weight: 600;">{sku_val}</span>
        </div>
        
        <div style="display: flex; gap: 20px; flex-wrap: wrap;">
            <div style="flex: 1.2; min-width: 280px; display: flex; flex-direction: column; gap: 15px;">
                <div>
                    <h4 style="margin: 0 0 6px 0; color: #8B6F4E; font-size: 13px; text-transform: uppercase; letter-spacing: 0.5px; border-bottom: 1px solid #E4D6CA; padding-bottom: 3px;">Details</h4>
                    {details_html or '<div style="color:#C4B0A0; font-size:12px; font-style:italic;">No metadata mapped</div>'}
                </div>
                
                <div style="display: flex; gap: 15px;">
                    <div style="flex: 1;">
                        <h4 style="margin: 0 0 6px 0; color: #8B6F4E; font-size: 12px; text-transform: uppercase; letter-spacing: 0.5px;">Key Features</h4>
                        <ul style="margin: 0; padding-left: 16px; color: #5C483A; font-size: 11.5px; line-height: 1.4;">
                            {features_html}
                        </ul>
                    </div>
                    <div style="flex: 1;">
                        <h4 style="margin: 0 0 6px 0; color: #8B6F4E; font-size: 12px; text-transform: uppercase; letter-spacing: 0.5px;">Benefits</h4>
                        <ul style="margin: 0; padding-left: 16px; color: #5C483A; font-size: 11.5px; line-height: 1.4;">
                            {benefits_html}
                        </ul>
                    </div>
                </div>
            </div>
            
            <div style="flex: 0.8; min-width: 220px; display: flex; align-items: center; justify-content: center;">
                {img_html}
            </div>
        </div>
    </div>
    """
    components.html(
        f"""<!DOCTYPE html><html><head><link href='https://fonts.googleapis.com/css2?family=DM+Sans:wght@300;400;500;600&display=swap' rel='stylesheet'></head><body style='margin:0;padding:0;background:transparent;'>{html_content}</body></html>""",
        height=420,
        scrolling=False
    )

def run_automation(excel_bytes, pptx_bytes, from_row, to_row, mapping_dict, image_mappings):
    if image_mappings is None:
        image_mappings = {}
    df = pd.read_excel(io.BytesIO(excel_bytes))

    start_idx = max(0, from_row - 2)
    end_idx = min(to_row - 1, len(df))
    df_subset = df.iloc[start_idx:end_idx].copy()

    if df_subset.empty:
        raise ValueError("Selected row range contains no data.")

    col_to_images = get_excel_images(excel_bytes)
    df_cols = list(df.columns)

    prs = Presentation(io.BytesIO(pptx_bytes))
    if len(prs.slides) != 1:
        raise ValueError("Template must have exactly 1 slide.")

    source_slide = prs.slides[0]
    for _ in range(len(df_subset) - 1):
        clone_slide(prs, source_slide)

    for i, (original_idx, row) in enumerate(df_subset.iterrows()):
        slide = prs.slides[i]

        replacements = {}
        for tag, cfg in mapping_dict.items():
            col    = cfg.get("column", "")
            fmt    = cfg.get("format", "Text")
            symbol = cfg.get("symbol", "$")
            if not col:
                continue
            val = row.get(col, "")
            if pd.isna(val):
                val = ""
            if fmt == "Currency":
                formatted_val = safe_format(val, is_currency=True, currency_symbol=symbol)
            elif fmt == "Percentage":
                formatted_val = safe_format(val, is_percent=True)
            elif fmt == "Integer":
                formatted_val = safe_format(val)
            else:
                formatted_val = safe_text(val)
            replacements[tag] = formatted_val

        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)

        for shape in slide.shapes:
            purge_empty_paragraphs(shape)

        if image_mappings:
            for img_tag, col_name in image_mappings.items():
                if not col_name or col_name not in df_cols:
                    continue
                col_idx = df_cols.index(col_name)
                imgs = col_to_images.get(col_idx, [])
                img_bytes = imgs[i] if i < len(imgs) else None

                ph_shape = None
                for shape in slide.shapes:
                    if shape.has_text_frame and img_tag in shape.text_frame.text:
                        ph_shape = shape
                        break
                if ph_shape:
                    if img_bytes:
                        insert_image_into_placeholder(slide, img_bytes, ph_shape)
                    else:
                        for para in ph_shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ""
                        ph_shape.fill.background()
        else:
            placeholder = find_image_placeholder(slide)
            if placeholder:
                all_imgs = [b for imgs in col_to_images.values() for b in imgs]
                img_bytes = all_imgs[i] if i < len(all_imgs) else None
                if img_bytes:
                    insert_image_into_placeholder(slide, img_bytes, placeholder)
                else:
                    for para in placeholder.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ""
                    placeholder.fill.background()

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue(), len(df_subset)


# ══════════════════════════════════════════════════════════════
#  TEMPLATE & EXCEL LOADERS
# ══════════════════════════════════════════════════════════════

TEMPLATES_DIR       = "templates"
EXCEL_TEMPLATES_DIR = "excel_templates"

def get_available_templates():
    if not os.path.exists(TEMPLATES_DIR):
        return []
    return [f for f in os.listdir(TEMPLATES_DIR) if f.endswith('.pptx')]

def load_template(filename):
    with open(os.path.join(TEMPLATES_DIR, filename), 'rb') as f:
        return f.read()

def get_available_excel_templates():
    if not os.path.exists(EXCEL_TEMPLATES_DIR):
        return []
    return [f for f in os.listdir(EXCEL_TEMPLATES_DIR) if f.endswith('.xlsx')]

def load_excel_template(filename):
    with open(os.path.join(EXCEL_TEMPLATES_DIR, filename), 'rb') as f:
        return f.read()


# ══════════════════════════════════════════════════════════════
#  HISTORY DATABASE UTILITIES
# ══════════════════════════════════════════════════════════════

HISTORY_FILE = "history.json"

def load_history():
    if os.path.exists(HISTORY_FILE):
        try:
            with open(HISTORY_FILE, "r") as f:
                return json.load(f)
        except:
            pass
    return []

def add_to_history(filename, template_name, excel_name, row_count, status="Success"):
    import datetime
    history = load_history()
    entry = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "filename": filename,
        "template": template_name,
        "excel": excel_name,
        "rows": row_count,
        "status": status
    }
    history.insert(0, entry)
    try:
        with open(HISTORY_FILE, "w") as f:
            json.dump(history, f, indent=4)
    except:
        pass


# ══════════════════════════════════════════════════════════════
#  WIZARD UI HELPERS
# ══════════════════════════════════════════════════════════════

def render_conversion_status(excel_file, mapping_dict, generating=False, ready=False):
    is_uploaded = excel_file is not None
    if is_uploaded:
        excel_name = excel_file if isinstance(excel_file, str) else excel_file.name
    else:
        excel_name = "Upload an Excel file"
        
    s1_class = "done" if is_uploaded else "todo"
    s1_icon = "✓" if is_uploaded else "1"
    s1_sub = excel_name
    
    s2_class = "done" if (is_uploaded and len(mapping_dict) > 0) else "todo"
    s2_icon = "✓" if (is_uploaded and len(mapping_dict) > 0) else "2"
    s2_sub = f"{len(mapping_dict)} fields mapped" if (is_uploaded and len(mapping_dict) > 0) else "Waiting for columns"
    
    s3_class = "active" if generating else ("done" if ready else "todo")
    s3_icon = "⏳" if generating else ("✓" if ready else "3")
    s3_sub = "Processing slides..." if generating else ("Slides generated!" if ready else "Waiting to start")
    
    s4_class = "done" if ready else "todo"
    s4_icon = "✓" if ready else "4"
    s4_sub = "Ready for download" if ready else "Waiting to start"
    
    html = f"""
    <div class="timeline-container" style="border:none; padding:0; box-shadow:none; background:transparent; margin-top:5px;">
        <div class="timeline-title" style="font-size:12px; font-weight:700; color:#2C1F14; text-transform:uppercase; margin-bottom:12px; letter-spacing:0.5px;">Conversion Status</div>
        <div class="timeline-item {s1_class}">
            <div class="timeline-badge">{s1_icon}</div>
            <div class="timeline-content">
                <div class="timeline-label">Excel file uploaded</div>
                <div class="timeline-desc">{s1_sub}</div>
            </div>
        </div>
        <div class="timeline-item {s2_class}">
            <div class="timeline-badge">{s2_icon}</div>
            <div class="timeline-content">
                <div class="timeline-label">Data validated</div>
                <div class="timeline-desc">{s2_sub}</div>
            </div>
        </div>
        <div class="timeline-item {s3_class}">
            <div class="timeline-badge">{s3_icon}</div>
            <div class="timeline-content">
                <div class="timeline-label">Generating slides</div>
                <div class="timeline-desc">{s3_sub}</div>
            </div>
        </div>
        <div class="timeline-item {s4_class}">
            <div class="timeline-badge">{s4_icon}</div>
            <div class="timeline-content">
                <div class="timeline-label">Ready to download</div>
                <div class="timeline-desc">{s4_sub}</div>
            </div>
        </div>
    </div>
    """
    st.markdown(html, unsafe_allow_html=True)


def step_card_open(num, title, desc=""):
    # Kept for legacy compatibility
    pass

def step_card_close():
    # Kept for legacy compatibility
    pass

def get_pptx_bytes():
    if "start_option" not in st.session_state:
        return None
    if st.session_state.start_option == "custom":
        return st.session_state.custom_template_bytes
    elif st.session_state.start_option == "pre_existing" and st.session_state.selected_template:
        try:
            return load_template(st.session_state.selected_template)
        except:
            return None
    return None

def reset_conversion_results():
    st.session_state.generating_state = False
    st.session_state.ready_state = False
    st.session_state.result_bytes = None
    st.session_state.count_slides = 0
    st.session_state.elapsed_time = 0.0
    st.session_state.error_msg = None
    st.session_state.output_filename = ""


# ══════════════════════════════════════════════════════════════
#  SIDEBAR NAVIGATION & STATE MANAGEMENT
# ══════════════════════════════════════════════════════════════

# Initialize progressive flow wizard state
if "start_option" not in st.session_state:
    st.session_state.start_option = None
if "selected_template" not in st.session_state:
    st.session_state.selected_template = None
if "custom_template_bytes" not in st.session_state:
    st.session_state.custom_template_bytes = None
if "excel_file_uploaded" not in st.session_state:
    st.session_state.excel_file_uploaded = False
if "excel_file_name" not in st.session_state:
    st.session_state.excel_file_name = None
if "excel_file_bytes" not in st.session_state:
    st.session_state.excel_file_bytes = None

if "generating_state" not in st.session_state:
    st.session_state.generating_state = False
if "ready_state" not in st.session_state:
    st.session_state.ready_state = False
if "result_bytes" not in st.session_state:
    st.session_state.result_bytes = None
if "count_slides" not in st.session_state:
    st.session_state.count_slides = 0
if "elapsed_time" not in st.session_state:
    st.session_state.elapsed_time = 0.0
if "error_msg" not in st.session_state:
    st.session_state.error_msg = None
if "output_filename" not in st.session_state:
    st.session_state.output_filename = ""

# Handle current page state
if "current_page" not in st.session_state:
    try:
        url_page = st.query_params.get("page", "create")
    except:
        url_page = "create"
    st.session_state.current_page = url_page if url_page in ["create", "history", "gallery"] else "create"

current_page = st.session_state.current_page

# Load Logo base64
logo_b64 = ""
logo_path = os.path.join(project_dir, "msi_logo.png")
if os.path.exists(logo_path):
    try:
        with open(logo_path, "rb") as f:
            logo_b64 = base64.b64encode(f.read()).decode("utf-8")
    except:
        pass

# Sidebar Logo Header
sidebar_logo_html = f"""
<div class="logo-box" style="text-align: center; padding: 20px 10px; border-bottom: 1px solid #EDE8E1; margin-bottom: 20px;">
    {"<img class='logo-img' src='data:image/png;base64," + logo_b64 + "' style='max-height:55px; max-width:100%; object-fit:contain; margin-bottom:8px; display:block; margin-left:auto; margin-right:auto;'/>" if logo_b64 else "<div class='logo-text' style='font-size:20px; font-weight:800; color:#8B6F4E; letter-spacing:2px;'>MSI SERVICES</div>"}
    <div class="logo-sub" style="font-size:10.5px; color:#8B6F4E; font-weight:600; letter-spacing:0.2px; line-height:1.3; opacity:0.95; margin-top:4px;">Making Dream Surfaces Attainable</div>
</div>
"""
st.sidebar.markdown(sidebar_logo_html, unsafe_allow_html=True)

# Sidebar Menu (using native buttons styled with CSS)
st.sidebar.markdown("<div class='menu-label'>CONVERSION</div>", unsafe_allow_html=True)
if st.sidebar.button("📄 Create Presentation", key="btn_nav_create", use_container_width=True, type="primary" if current_page == "create" else "secondary"):
    st.session_state.current_page = "create"
    st.rerun()

if st.sidebar.button("📜 History Log", key="btn_nav_history", use_container_width=True, type="primary" if current_page == "history" else "secondary"):
    st.session_state.current_page = "history"
    st.rerun()

st.sidebar.markdown("<div class='menu-label' style='margin-top:15px;'>TEMPLATES</div>", unsafe_allow_html=True)
if st.sidebar.button("🎨 Template Gallery", key="btn_nav_gallery", use_container_width=True, type="primary" if current_page == "gallery" else "secondary"):
    st.session_state.current_page = "gallery"
    st.rerun()


# ══════════════════════════════════════════════════════════════
#  HEADER BAR
# ══════════════════════════════════════════════════════════════

col_title, col_help = st.columns([4, 1])
with col_title:
    if current_page == "create":
        st.markdown("<h2 style='margin:0;color:#2C1F14;'>Create Presentation</h2><p style='margin:4px 0 0;font-size:13px;color:#9A8070;'>Convert your Excel data into stunning PowerPoint slides in seconds.</p>", unsafe_allow_html=True)
    elif current_page == "history":
        st.markdown("<h2 style='margin:0;color:#2C1F14;'>History</h2><p style='margin:4px 0 0;font-size:13px;color:#9A8070;'>Log of previously generated slide presentations.</p>", unsafe_allow_html=True)
    elif current_page == "gallery":
        st.markdown("<h2 style='margin:0;color:#2C1F14;'>Template Gallery</h2><p style='margin:4px 0 0;font-size:13px;color:#9A8070;'>Manage templates and upload custom ones for regular use.</p>", unsafe_allow_html=True)

with col_help:
    help_popover = st.popover("❓ Help Support", use_container_width=True)
    with help_popover:
        st.markdown("""
        <div style="font-family:'Plus Jakarta Sans',sans-serif;padding:4px;">
            <h4 style="margin:0 0 10px;color:#8B6F4E;font-size:13px;border-bottom:1px solid #EDE8E1;padding-bottom:6px;">Support Contact</h4>
            <p style="margin:4px 0;font-size:11.5px;color:#2C1F14;"><b>Name:</b> Sarthak Sharma</p>
            <p style="margin:4px 0;font-size:11.5px;color:#2C1F14;"><b>Email:</b> <a href="mailto:sarthak.s@msisurfaces.com" style="color:#8B6F4E;text-decoration:none;">sarthak.s@msisurfaces.com</a></p>
            <p style="margin:4px 0;font-size:11.5px;color:#2C1F14;"><b>Ext:</b> 5023</p>
            <hr style="border:none;border-top:1px solid #EDE8E1;margin:8px 0;"/>
            <a href="https://teams.microsoft.com/l/chat/0/0?users=sarthak.s@msisurfaces.com" target="_blank" style="display:block;text-align:center;background:#8B6F4E;color:#fff;padding:8px 12px;border-radius:6px;font-size:11px;font-weight:700;text-decoration:none;transition:background 0.2s;">💬 Chat on Teams</a>
        </div>
        """, unsafe_allow_html=True)

st.markdown("<div class='fancy-divider'></div>", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
#  PAGE WORKSPACE ROUTER
# ══════════════════════════════════════════════════════════════

if current_page == "create":
    available_templates = get_available_templates()

    # Flow Reset header button
    col_flow, col_reset = st.columns([5, 1])
    with col_flow:
        pass
    with col_reset:
        if st.button("🔄 Reset Flow", key="btn_reset_flow", use_container_width=True):
            st.session_state.start_option = None
            st.session_state.selected_template = None
            st.session_state.custom_template_bytes = None
            st.session_state.excel_file_uploaded = False
            st.session_state.excel_file_name = None
            st.session_state.excel_file_bytes = None
            reset_conversion_results()
            st.rerun()

    col_left, col_right = st.columns([13, 7])

    # Left workspace progressive wizard steps
    with col_left:
        # Step 1: Start selection card
        with st.container(border=True):
            st.markdown("""
            <div style="display:flex; align-items:center; gap:12px; margin-bottom:15px;">
                <div class="step-card-num">1</div>
                <div>
                    <div class="step-card-title" style="font-size:15px; font-weight:700; color:#2C1F14;">How would you like to start?</div>
                    <div class="step-card-desc" style="font-size:11.5px; color:#9A8070;">Choose the template source for your presentation slides</div>
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            c1, c2 = st.columns(2)
            with c1:
                is_pre = (st.session_state.start_option == "pre_existing")
                if st.button(
                    "📁 Use Pre-existing Template", 
                    key="btn_start_pre_existing", 
                    use_container_width=True, 
                    type="primary" if is_pre else "secondary"
                ):
                    st.session_state.start_option = "pre_existing"
                    st.session_state.selected_template = None
                    st.session_state.custom_template_bytes = None
                    st.session_state.excel_file_uploaded = False
                    st.session_state.excel_file_name = None
                    st.session_state.excel_file_bytes = None
                    reset_conversion_results()
                    st.rerun()
            with c2:
                is_cust = (st.session_state.start_option == "custom")
                if st.button(
                    "📤 Upload Custom Template", 
                    key="btn_start_custom", 
                    use_container_width=True, 
                    type="primary" if is_cust else "secondary"
                ):
                    st.session_state.start_option = "custom"
                    st.session_state.selected_template = None
                    st.session_state.custom_template_bytes = None
                    st.session_state.excel_file_uploaded = False
                    st.session_state.excel_file_name = None
                    st.session_state.excel_file_bytes = None
                    reset_conversion_results()
                    st.rerun()

        # Step 2: Choose / Upload template (progressive)
        if st.session_state.start_option is not None:
            with st.container(border=True):
                if st.session_state.start_option == "pre_existing":
                    st.markdown("""
                    <div style="display:flex; align-items:center; gap:12px; margin-bottom:15px;">
                        <div class="step-card-num">2</div>
                        <div>
                            <div class="step-card-title" style="font-size:15px; font-weight:700; color:#2C1F14;">Select Slide Template</div>
                            <div class="step-card-desc" style="font-size:11.5px; color:#9A8070;">Choose from our catalog of pre-uploaded templates</div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    if not available_templates:
                        st.warning("📂 No templates found in the `templates/` folder. Upload templates in the Template Gallery page.")
                    else:
                        selected = st.selectbox(
                            "Select Template",
                            options=available_templates,
                            index=None if st.session_state.selected_template is None else available_templates.index(st.session_state.selected_template) if st.session_state.selected_template in available_templates else None,
                            placeholder="Select a template from catalog...",
                            format_func=lambda x: x.replace('.pptx', '').replace('_', ' '),
                            key="sel_pre_template"
                        )
                        if selected != st.session_state.selected_template:
                            st.session_state.selected_template = selected
                            st.session_state.excel_file_uploaded = False
                            st.session_state.excel_file_name = None
                            st.session_state.excel_file_bytes = None
                            reset_conversion_results()
                            st.rerun()
                
                else: # start_option == "custom"
                    st.markdown("""
                    <div style="display:flex; align-items:center; gap:12px; margin-bottom:15px;">
                        <div class="step-card-num">2</div>
                        <div>
                            <div class="step-card-title" style="font-size:15px; font-weight:700; color:#2C1F14;">Upload Custom Template</div>
                            <div class="step-card-desc" style="font-size:11.5px; color:#9A8070;">Upload a single-slide .pptx presentation file containing placeholders</div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    custom_pptx = st.file_uploader(
                        "Upload PowerPoint Template (.pptx)",
                        type=["pptx"],
                        help="Upload a single-slide .pptx with [TAG (Type)] placeholders.",
                        key="custom_pptx_uploader",
                        label_visibility="collapsed"
                    )
                    if custom_pptx:
                        if st.session_state.selected_template != custom_pptx.name or st.session_state.custom_template_bytes is None:
                            st.session_state.selected_template = custom_pptx.name
                            st.session_state.custom_template_bytes = custom_pptx.getvalue()
                            st.session_state.excel_file_uploaded = False
                            st.session_state.excel_file_name = None
                            st.session_state.excel_file_bytes = None
                            reset_conversion_results()
                            st.rerun()
                    else:
                        if st.session_state.selected_template is not None:
                            st.session_state.selected_template = None
                            st.session_state.custom_template_bytes = None
                            reset_conversion_results()
                            st.rerun()

        # Step 3: Upload Excel product data (progressive)
        if st.session_state.selected_template is not None:
            with st.container(border=True):
                st.markdown("""
                <div style="display:flex; align-items:center; gap:12px; margin-bottom:15px;">
                    <div class="step-card-num">3</div>
                    <div>
                        <div class="step-card-title" style="font-size:15px; font-weight:700; color:#2C1F14;">Upload Product Data</div>
                        <div class="step-card-desc" style="font-size:11.5px; color:#9A8070;">Excel sheet with column details and embedded product images</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
                
                # Starter Excel file downloader
                excel_templates = get_available_excel_templates()
                if excel_templates:
                    st.markdown("<p style='font-size:11px;font-weight:700;color:#8B6F4E;margin-bottom:6px;text-transform:uppercase;'>⬇️ Starter Excel Templates</p>", unsafe_allow_html=True)
                    _tc1, _tc2 = st.columns([3, 1])
                    with _tc1:
                        selected_excel = st.selectbox(
                            "Excel Template",
                            options=excel_templates,
                            index=None,
                            placeholder="Select starter template to download...",
                            format_func=lambda x: x.replace('.xlsx','').replace('_',' '),
                            label_visibility="collapsed",
                            key="sel_starter_excel"
                        )
                    with _tc2:
                        if selected_excel:
                            st.download_button(
                                label="📥 Get Template",
                                data=load_excel_template(selected_excel),
                                file_name=selected_excel,
                                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                key="dl_excel_template",
                                use_container_width=True
                            )
                    st.markdown("<div style='height:12px;'></div>", unsafe_allow_html=True)

                excel_file = st.file_uploader(
                    "Upload Excel File (.xlsx)",
                    type=["xlsx"],
                    help="Excel sheet with images. Columns match template tags.",
                    label_visibility="collapsed",
                    key="excel_data_uploader"
                )
                
                if excel_file:
                    if not st.session_state.excel_file_uploaded or st.session_state.excel_file_name != excel_file.name:
                        st.session_state.excel_file_uploaded = True
                        st.session_state.excel_file_name = excel_file.name
                        st.session_state.excel_file_bytes = excel_file.getvalue()
                        reset_conversion_results()
                        st.rerun()
                        
                    st.markdown(f"""
                    <div style="background:#FFF8F2; border:1px solid #E8D9CA; border-radius:10px; padding:12px 16px; display:flex; align-items:center; justify-content:space-between; margin-top:10px;">
                        <div style="display:flex; align-items:center; gap:12px;">
                            <span style="font-size:24px;">📊</span>
                            <div>
                                <div style="font-weight:700; color:#2C1F14; font-size:13px;">{excel_file.name}</div>
                                <div style="font-size:11px; color:#9A8070;">{excel_file.size / 1024:.1f} KB &nbsp;·&nbsp; <span style="color:#0F5132; font-weight:600;">✔️ File uploaded successfully</span></div>
                            </div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    if st.session_state.excel_file_uploaded:
                        st.session_state.excel_file_uploaded = False
                        st.session_state.excel_file_name = None
                        st.session_state.excel_file_bytes = None
                        reset_conversion_results()
                        st.rerun()

        # Step 4: Configure & Convert (progressive)
        mapping_dict = {}
        image_mappings = {}
        if st.session_state.excel_file_uploaded:
            with st.container(border=True):
                st.markdown("""
                <div style="display:flex; align-items:center; gap:12px; margin-bottom:15px;">
                    <div class="step-card-num">4</div>
                    <div>
                        <div class="step-card-title" style="font-size:15px; font-weight:700; color:#2C1F14;">Configure Data & Run</div>
                        <div class="step-card-desc" style="font-size:11.5px; color:#9A8070;">Customize range and filename, then generate PowerPoint</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
                
                _s3c1, _s3c2 = st.columns([2, 1])
                with _s3c1:
                    output_name = st.text_input(
                        "Output Filename",
                        value="",
                        placeholder="e.g. Sales_Product_Proposal",
                        help=".pptx extension is added automatically",
                        key="txt_output_name"
                    )
                with _s3c2:
                    st.markdown("&nbsp;", unsafe_allow_html=True)
                    
                _row_c1, _row_c2 = st.columns(2)
                with _row_c1:
                    from_row = st.number_input(
                        "From Row",
                        min_value=1, max_value=9999, value=2,
                        help="First data row to process (row 2 = first product)",
                        key="num_from_row"
                    )
                with _row_c2:
                    to_row = st.number_input(
                        "To Row",
                        min_value=1, max_value=9999, value=9999,
                        help="Last data row to process (9999 = all rows)",
                        key="num_to_row"
                    )
                
                st.markdown("<hr style='border:none;border-top:1px solid #EDE8E1;margin:15px 0;'/>", unsafe_allow_html=True)
                
                # Check mapping first
                pptx_bytes = get_pptx_bytes()
                excel_bytes = st.session_state.excel_file_bytes
                
                if excel_bytes and pptx_bytes:
                    try:
                        _detected_tags = extract_placeholders_from_pptx(pptx_bytes)
                        _df_am = pd.read_excel(io.BytesIO(excel_bytes))
                        mapping_dict, image_mappings = build_auto_mapping(_detected_tags, list(_df_am.columns))
                    except:
                        pass
                        
                generate_clicked = st.button("🚀 Start Conversion", key="btn_generate")

                if generate_clicked:
                    if not excel_bytes:
                        st.error("Please upload an Excel data file first.")
                    elif pptx_bytes is None:
                        st.error("Please select or upload a template first.")
                    else:
                        st.session_state.generating_state = True
                        st.session_state.ready_state = False
                        st.session_state.error_msg = None
                        
                        _tmpl_label = (st.session_state.selected_template or "").replace('.pptx','').replace('_',' ')
                        default_fname = f"{(_tmpl_label or 'Output').replace(' ','_')}_Output"
                        fname = output_name.strip() or default_fname
                        if not fname.endswith('.pptx'):
                            fname += '.pptx'
                        st.session_state.output_filename = fname
                        
                        import time as _time
                        _t0 = _time.time()
                        
                        try:
                            result_bytes, count_slides = run_automation(
                                excel_bytes, pptx_bytes,
                                from_row=int(from_row),
                                to_row=int(to_row),
                                mapping_dict=mapping_dict,
                                image_mappings=image_mappings
                            )
                            _elapsed = round(_time.time() - _t0, 1)
                            
                            st.session_state.result_bytes = result_bytes
                            st.session_state.count_slides = count_slides
                            st.session_state.elapsed_time = _elapsed
                            st.session_state.ready_state = True
                            st.session_state.generating_state = False
                            
                            add_to_history(
                                filename=fname,
                                template_name=st.session_state.selected_template or "Custom Upload",
                                excel_name=st.session_state.excel_file_name,
                                row_count=count_slides,
                                status="Success"
                            )
                            st.rerun()
                        except Exception as e:
                            st.session_state.ready_state = False
                            st.session_state.generating_state = False
                            st.session_state.error_msg = str(e)
                            add_to_history(
                                filename=fname,
                                template_name=st.session_state.selected_template or "Custom Upload",
                                excel_name=st.session_state.excel_file_name,
                                row_count=0,
                                status="Failed"
                            )
                            st.error(f"Generation error: {e}")
                            st.rerun()

    # Right workspace column (Preview, timeline, tip merged in single block)
    with col_right:
        with st.container(border=True):
            col_r_lbl, col_r_btn = st.columns([5, 4])
            with col_r_lbl:
                st.markdown("<h4 style='margin:8px 0 0;color:#2C1F14;font-size:14px;font-weight:700;'>Workspace Overview</h4>", unsafe_allow_html=True)
            with col_r_btn:
                if st.button("📜 View History", key="btn_go_history_page", use_container_width=True):
                    st.session_state.current_page = "history"
                    st.rerun()
            
            st.markdown("<div style='height:12px; border-bottom:1px solid #EDE8E1; margin-bottom:12px;'></div>", unsafe_allow_html=True)

            # Template Preview
            st.markdown("<div class='timeline-title' style='margin-bottom:8px;font-size:12px;font-weight:700;color:#2C1F14;text-transform:uppercase;'>Template Preview</div>", unsafe_allow_html=True)
            
            pptx_bytes = get_pptx_bytes()
            if pptx_bytes is not None and len(mapping_dict) > 0:
                if st.session_state.excel_file_uploaded and st.session_state.excel_file_bytes:
                    try:
                        df = pd.read_excel(io.BytesIO(st.session_state.excel_file_bytes))
                        row_idx = max(0, int(from_row) - 2)
                        if row_idx < len(df):
                            render_slide_preview(
                                mapping_dict, 
                                image_mappings=image_mappings, 
                                excel_row=df.iloc[row_idx].to_dict(),
                                excel_row_idx=row_idx,
                                is_template_mode=False
                            )
                        else:
                            st.info("Row index out of range for preview.")
                    except Exception as e:
                        st.info(f"Unable to parse Excel preview: {e}")
                else:
                    render_slide_preview(mapping_dict, is_template_mode=True)
            else:
                st.markdown("""
                <div style="background:#FAF8F5;border:1px dashed #CFC0B0;border-radius:8px;height:180px;display:flex;align-items:center;justify-content:center;color:#B8A898;font-size:12px;font-style:italic;margin-bottom:15px;text-align:center;padding:15px;">
                    Select/upload a template and Excel data to see live slide preview here
                </div>
                """, unsafe_allow_html=True)

            st.markdown("<div style='height:12px; border-bottom:1px solid #EDE8E1; margin-bottom:12px;'></div>", unsafe_allow_html=True)

            # Result status card (success or failure)
            if st.session_state.ready_state and st.session_state.result_bytes:
                st.markdown(f"""
                <div class="result-card" style="margin-bottom:15px;">
                  <div class="result-icon">🎉</div>
                  <div class="result-title">Presentation Ready!</div>
                  <div class="result-sub">{st.session_state.output_filename}</div>
                  <div class="result-stats" style="margin-top:10px; margin-bottom:15px;">
                    <div class="stat-chip">
                      <span class="val">{st.session_state.count_slides}</span>
                      <span class="lbl">Slide{'s' if st.session_state.count_slides!=1 else ''}</span>
                    </div>
                    <div class="stat-chip">
                      <span class="val">{len(mapping_dict)}</span>
                      <span class="lbl">Fields</span>
                    </div>
                    <div class="stat-chip">
                      <span class="val">{st.session_state.elapsed_time}s</span>
                      <span class="lbl">Time</span>
                    </div>
                  </div>
                </div>
                """, unsafe_allow_html=True)

                st.download_button(
                    label=f"📥 Download {st.session_state.output_filename}",
                    data=st.session_state.result_bytes,
                    file_name=st.session_state.output_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_result_sidebar",
                    use_container_width=True
                )
                st.markdown("<div style='height:12px;'></div>", unsafe_allow_html=True)

            elif st.session_state.error_msg:
                st.error(f"Generation error: {st.session_state.error_msg}")

            # Conversion Status Timeline
            excel_name_for_timeline = st.session_state.excel_file_name if st.session_state.excel_file_uploaded else None
            render_conversion_status(
                excel_file=excel_name_for_timeline, 
                mapping_dict=mapping_dict, 
                generating=st.session_state.generating_state, 
                ready=st.session_state.ready_state
            )
            
            st.markdown("<div style='height:12px; border-bottom:1px solid #EDE8E1; margin-bottom:12px;'></div>", unsafe_allow_html=True)

            # Tip Card
            st.markdown("""
            <div style="background:#FFF8F2; border:1px solid #E8D9CA; border-radius:8px; padding:12px; display:flex; gap:10px; align-items:flex-start;">
                <span style="font-size:16px;">💡</span>
                <div>
                    <div style="font-weight:700; color:#8B6F4E; font-size:12px; margin-bottom:2px;">Placeholder Tip</div>
                    <div style="font-size:11px; color:#5C483A; line-height:1.4;">
                        Ensure the placeholders in your template are named exactly like the Excel column they replace, followed by their format in parentheses, all enclosed in square brackets. For example: <code>[RETAIL (Currency $)]</code>, <code>[IMAGE 1 (Image)]</code>, or <code>[NAME (Text)]</code>.
                    </div>
                </div>
            </div>
            """, unsafe_allow_html=True)


elif current_page == "history":
    with st.container(border=True):
        st.markdown("<h4 style='margin:0 0 8px 0;color:#2C1F14;font-size:14px;font-weight:700;'>Conversion Run History</h4>", unsafe_allow_html=True)
        st.markdown("<p style='font-size:12px;color:#9A8070;margin-bottom:20px;'>Summary of generated files (zero local storage used):</p>", unsafe_allow_html=True)
        
        history_list = load_history()
        if not history_list:
            st.info("No conversion history found.")
        else:
            table_rows = ""
            for idx, entry in enumerate(history_list):
                status_color = "#385723" if entry.get('status') == "Success" else "#C00000"
                status_bg = "#E2F0D9" if entry.get('status') == "Success" else "#FCE4D6"
                table_rows += f"""
                <tr style="border-bottom:1px solid #EDE8E1;">
                    <td style="padding:12px;font-size:12px;color:#2C1F14;font-weight:600;">{entry.get('timestamp')}</td>
                    <td style="padding:12px;font-size:12px;color:#8B6F4E;font-weight:700;">{entry.get('filename')}</td>
                    <td style="padding:12px;font-size:12px;color:#5C483A;">{entry.get('template')}</td>
                    <td style="padding:12px;font-size:12px;color:#5C483A;">{entry.get('excel')}</td>
                    <td style="padding:12px;font-size:12px;color:#2C1F14;font-weight:600;text-align:center;">{entry.get('rows')}</td>
                    <td style="padding:12px;font-size:11px;text-align:center;">
                        <span style="background:{status_bg};color:{status_color};padding:2px 8px;border-radius:20px;font-weight:600;">{entry.get('status', 'Success')}</span>
                    </td>
                </tr>
                """
            
            st.markdown(f"""
            <div style="overflow-x:auto;">
                <table style="width:100%;border-collapse:collapse;text-align:left;">
                    <thead>
                        <tr style="background:#FAF8F5;border-bottom:2px solid #EDE8E1;">
                            <th style="padding:12px;font-size:11px;font-weight:700;color:#9A8070;text-transform:uppercase;">Date & Time</th>
                            <th style="padding:12px;font-size:11px;font-weight:700;color:#9A8070;text-transform:uppercase;">Output Name</th>
                            <th style="padding:12px;font-size:11px;font-weight:700;color:#9A8070;text-transform:uppercase;">Template</th>
                            <th style="padding:12px;font-size:11px;font-weight:700;color:#9A8070;text-transform:uppercase;">Excel File</th>
                            <th style="padding:12px;font-size:11px;font-weight:700;color:#9A8070;text-transform:uppercase;text-align:center;">Rows</th>
                            <th style="padding:12px;font-size:11px;font-weight:700;color:#9A8070;text-transform:uppercase;text-align:center;">Status</th>
                        </tr>
                    </thead>
                    <tbody>
                        {table_rows}
                    </tbody>
                </table>
            </div>
            """, unsafe_allow_html=True)


elif current_page == "gallery":
    if not os.path.exists(TEMPLATES_DIR):
        try:
            os.makedirs(TEMPLATES_DIR)
        except:
            pass

    with st.container(border=True):
        st.markdown("""
        <div style="display:flex; align-items:center; gap:12px; margin-bottom:15px;">
            <div class="step-card-num" style="background:#8B6F4E;color:#fff;">➕</div>
            <div>
                <div class="step-card-title" style="font-size:15px; font-weight:700; color:#2C1F14;">Upload New Template</div>
                <div class="step-card-desc" style="font-size:11.5px; color:#9A8070;">Add a custom slide template (.pptx) to the permanent gallery</div>
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        new_template_file = st.file_uploader(
            "Upload PowerPoint Template (.pptx)",
            type=["pptx"],
            help="Upload a single-slide .pptx template with placeholders.",
            key="gallery_template_upload",
            label_visibility="collapsed"
        )
        if new_template_file:
            target_path = os.path.join(TEMPLATES_DIR, new_template_file.name)
            try:
                with open(target_path, "wb") as f:
                    f.write(new_template_file.getvalue())
                st.success(f"✔️ **{new_template_file.name}** has been uploaded and added to the gallery!")
                st.rerun()
            except Exception as e:
                st.error(f"Failed to save template: {e}")
                
    st.markdown("<div style='height:15px;'></div>", unsafe_allow_html=True)
    
    with st.container(border=True):
        st.markdown("<h4 style='color:#2C1F14;margin:0 0 15px 0;font-size:14px;font-weight:700;'>📁 Active Templates</h4>", unsafe_allow_html=True)
        
        templates = get_available_templates()
        if not templates:
            st.info("No templates available in the gallery. Upload custom templates above.")
        else:
            for t_name in templates:
                t_path = os.path.join(TEMPLATES_DIR, t_name)
                
                st.markdown(f"""
                <div style="background:#FAF8F5; border:1px solid #EDE8E1; border-radius:10px; padding:12px 16px; margin-bottom:12px;">
                    <div style="display:flex; align-items:center; justify-content:space-between;">
                        <div style="display:flex; align-items:center; gap:12px;">
                            <span style="font-size:24px;">📄</span>
                            <div>
                                <div style="font-weight:700; color:#2C1F14; font-size:13px;">{t_name.replace('.pptx','').replace('_',' ')}</div>
                                <div style="font-size:11px; color:#9A8070;">Filename: {t_name}</div>
                            </div>
                        </div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
                
                c_dl, c_del, c_spacer = st.columns([1.5, 1.5, 5])
                with c_dl:
                    try:
                        with open(t_path, "rb") as f:
                            st.download_button(
                                label="📥 Download",
                                data=f.read(),
                                file_name=t_name,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key=f"dl_{t_name}",
                                use_container_width=True
                            )
                    except:
                        pass
                with c_del:
                    if st.button("🗑️ Delete", key=f"del_{t_name}", use_container_width=True):
                        try:
                            os.remove(t_path)
                            st.success(f"Deleted {t_name}")
                            st.rerun()
                        except Exception as e:
                            st.error(f"Error: {e}")
                st.markdown("<hr style='border:none;border-top:1px solid #EDE8E1;margin:12px 0;'/>", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════
#  FOOTER
# ══════════════════════════════════════════════════════════════

st.markdown("""
<div class="msi-footer">
  MSI Services &nbsp;·&nbsp; Sales Operations Slide Automation Tool &nbsp;·&nbsp; Internal Support
</div>
""", unsafe_allow_html=True)
